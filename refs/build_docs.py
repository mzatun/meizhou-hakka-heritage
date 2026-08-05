# -*- coding: utf-8 -*-
"""从课程数据生成《梅州客家非遗精讲》Word 讲义 与 配套教学 PPT。
依赖：python-docx, python-pptx（已安装于托管版 Python）。
用法：python refs/build_docs.py
"""
import json, os
BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
def load(name):
    with open(os.path.join(BASE, "data", name + ".json"), encoding="utf-8") as f:
        return json.load(f)

chapters = load("chapters")["chapters"]
cities   = load("cities")
relics   = load("relics")["relics"]

INDIGO = "1F3A5F"; EARTH = "C9973F"; FLAME = "C0392B"; INK = "1A1A1A"

# ---------------- Word 讲义 ----------------
from docx import Document
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn

def set_cjk(run, name="Microsoft YaHei"):
    run.font.name = name
    rpr = run._element.get_or_add_rPr()
    rfonts = rpr.find(qn('w:rFonts'))
    if rfonts is None:
        from docx.oxml import OxmlElement
        rfonts = OxmlElement('w:rFonts')
        rpr.append(rfonts)
    rfonts.set(qn('w:eastAsia'), name)

doc = Document()
# 默认正文字体
normal = doc.styles['Normal']
normal.font.name = 'Microsoft YaHei'
normal.font.size = Pt(11)
normal.element.rPr.rFonts.set(qn('w:eastAsia'), 'Microsoft YaHei')

def H(text, level=1, color=INDIGO):
    h = doc.add_heading(level=level)
    r = h.add_run(text); r.font.color.rgb = RGBColor.from_string(color)
    set_cjk(r, "Microsoft YaHei")
    return h

def P(text, bold=False, italic=False, size=11, color=INK, align=None, space=6):
    p = doc.add_paragraph()
    r = p.add_run(text); r.bold = bold; r.italic = italic
    r.font.size = Pt(size); r.font.color.rgb = RGBColor.from_string(color)
    set_cjk(r, "Microsoft YaHei")
    if align: p.alignment = align
    p.paragraph_format.space_after = Pt(space)
    return p

def bullet(text, color=INK):
    p = doc.add_paragraph(style='List Bullet')
    r = p.add_run(text); r.font.size = Pt(10.5); r.font.color.rgb = RGBColor.from_string(color)
    set_cjk(r, "Microsoft YaHei")
    return p

# 封面
t = doc.add_paragraph(); t.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = t.add_run("梅州客家非遗精讲"); r.bold = True; r.font.size = Pt(26); r.font.color.rgb = RGBColor.from_string(INDIGO)
set_cjk(r)
s = doc.add_paragraph(); s.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = s.add_run("高校通识课程 · 沉浸式数字课程门户 · 配套讲义"); r.font.size = Pt(13); r.font.color.rgb = RGBColor.from_string(EARTH)
set_cjk(r)
m = doc.add_paragraph(); m.alignment = WD_ALIGN_PARAGRAPH.CENTER
r = m.add_run("客从何处来 · 非遗如何活"); r.italic = True; r.font.size = Pt(12); r.font.color.rgb = RGBColor.from_string(FLAME)
set_cjk(r)
doc.add_paragraph()

# 导言
H("课程导言", 1)
P("《梅州客家非遗》是一门面向高校通识教育的 immersive 数字课程。课程以客家民系的千年迁徙为经，以梅州真实的非物质文化遗产代表性项目为纬，"
  "串联起山歌、汉剧、汉乐、埔寨火龙、围龙屋、客家饮食、礼俗与侨批过番八大主题。本课程配套数字门户（首页、课程章节、互动地图、数字非遗长廊、AI 讲解员、课程档案馆），"
  "倡导在真实可考据的非遗知识之上，用博物馆级的视觉与交互，让年轻一代“走进”而非仅仅“阅读”客家文化。")
P("本讲义依据各级非物质文化遗产名录与公开资料编写，图文视觉素材由 AI 生成，仅用于教学演示；非遗知识以官方公布信息为准。", italic=True, size=10, color=FLAME)

# 章节
for c in chapters:
    H(f"第{c['no']}章　{c['title']} · {c['subtitle']}", 2)
    meta = doc.add_paragraph()
    r = meta.add_run(f"时期：{c['period']}　|　级别：{c['level']}"); r.font.size = Pt(10); r.font.color.rgb = RGBColor.from_string(EARTH)
    set_cjk(r)
    P(c['summary'], bold=True, size=11.5)
    P("知识要点", bold=True, size=10.5, color=INDIGO, space=2)
    for p in c['points']:
        bullet(p)
    # 讲解稿
    sc = doc.add_paragraph()
    r = sc.add_run("【语音讲解稿】"); r.bold = True; r.font.size = Pt(10.5); r.font.color.rgb = RGBColor.from_string(FLAME)
    set_cjk(r)
    P(c['script'], size=10.5, color=INK)
    doc.add_paragraph()

# 附录：地图与长廊概览
H("附录一　客家迁徙与梅州非遗分布（互动地图）", 1)
P(f"互动地图标注 {sum(len(cities[k]) for k in ['route','counties','spread','overseas'])} 个节点，含迁徙路线、梅州八县市区非遗分布与海外播迁城市：")
for label, key in [("迁徙路线","route"),("梅州县区","counties"),("国内播迁","spread"),("海外播迁","overseas")]:
    names = "、".join(x.get('name','') for x in cities[key])
    bullet(f"{label}：{names}")

H("附录二　数字非遗长廊（40 件器物）", 1)
P("数字长廊以瀑布流展陈 40 件客家非遗器物，点击卡片可 360° 翻转，查看断代、工艺、流传区域、传承谱系与故事。举例如：")
for r0 in relics[:10]:
    bullet(f"{r0['name']}（{r0['era']}·{r0['region']}）—— {r0['craft']}")
P("……其余 30 件详见数字门户「非遗长廊」板块。", italic=True, size=10)

H("附录三　学习资源", 1)
bullet("课程数字门户：首页 / 课程章节 / 互动地图 / 数字非遗长廊 / 课程档案馆")
bullet("AI 讲解员「阿蛮」：右下角常驻，课程知识自由问答，支持流式呈现与语音播报")
bullet("配套教学 PPT：与本讲义同步发布，供课堂讲授使用")
bullet("延展阅读：罗香林《客家研究导论》、各级非物质文化遗产名录、梅州市非遗保护中心公开资料")

doc.save(os.path.join(BASE, "downloads", "meizhou-hakka-heritage-notes.docx"))
print("✓ docx 生成完成")

# ---------------- 教学 PPT ----------------
from pptx import Presentation
from pptx.util import Inches, Pt as PPT
from pptx.dml.color import RGBColor as PRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]
C_INDIGO = PRGB.from_string(INDIGO); C_EARTH = PRGB.from_string(EARTH); C_FLAME = PRGB.from_string(FLAME); C_WHITE = PRGB.from_string("FFFFFF"); C_INK = PRGB.from_string("1A1A1A")

def bg(slide, hexcolor):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PRGB.from_string(hexcolor)

def box(slide, l, t, w, h):
    return slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))

def setfont(run, size=18, color=C_INK, bold=False, name="Microsoft YaHei"):
    run.font.size = PPT(size); run.font.bold = bold; run.font.name = name
    run.font.color.rgb = color

def add(slide, l, t, w, h, lines, anchor=MSO_ANCHOR.TOP):
    tb = box(slide, l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, (txt, sz, col, bold) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run(); r.text = txt; setfont(r, sz, col, bold)
    return tb

# 封面
s = prs.slides.add_slide(blank); bg(s, INDIGO)
add(s, 1.2, 2.4, 10.9, 1.4, [("梅州客家非遗", 44, C_EARTH, True)])
add(s, 1.2, 3.7, 10.9, 1.0, [("高校通识课程 · 配套教学课件", 22, C_WHITE, False)])
add(s, 1.2, 5.6, 10.9, 0.8, [("客从何处来 · 非遗如何活", 18, C_FLAME, True)])

# 导览
s = prs.slides.add_slide(blank); bg(s, "F4EFE6")
add(s, 1.0, 0.7, 11.3, 1.0, [("课程导览", 36, C_INDIGO, True)])
for i, (big, lab) in enumerate([("8","课程章节"),("30","迁徙地图节点"),("40","数字非遗器物"),("4+","国家级非遗")]):
    x = 1.0 + i*3.05
    add(s, x, 2.3, 2.8, 1.6, [(big, 52, C_FLAME, True)])
    add(s, x, 3.9, 2.8, 0.8, [(lab, 16, C_INDIGO, False)])
add(s, 1.0, 5.4, 11.3, 1.4, [("从永嘉之乱的衣冠南渡，到今天活态传承的非遗技艺——以八章为经，以迁徙路线、三十节点、四十器物为纬，读懂世界客都的千年文脉。", 16, C_INK, False)])

# 每章
for c in chapters:
    s = prs.slides.add_slide(blank); bg(s, "FBF8F2")
    add(s, 0.9, 0.6, 11.5, 0.9, [(f"第{c['no']}章　{c['title']}", 32, C_INDIGO, True)])
    add(s, 0.9, 1.5, 11.5, 0.7, [(c['subtitle'] + "　·　" + c['level'], 16, C_EARTH, False)])
    add(s, 0.9, 2.4, 11.5, 0.9, [(c['summary'], 18, C_INK, False)])
    pts = c['points'][:4]
    lines = [(("• " + p), 15, C_INK, False) for p in pts]
    add(s, 0.9, 3.6, 11.5, 3.2, lines)
    add(s, 0.9, 6.9, 11.5, 0.6, [("时期：" + c['period'], 13, C_FLAME, True)])

# 地图概览
s = prs.slides.add_slide(blank); bg(s, INDIGO)
add(s, 1.0, 0.7, 11.3, 1.0, [("互动地图：客家迁徙与梅州非遗分布", 32, C_EARTH, True)])
add(s, 1.0, 2.0, 11.3, 0.8, [("30 个节点串起一千七百年的足迹", 18, C_WHITE, False)])
lines = []
for label, key in [("迁徙路线","route"),("梅州县区","counties"),("国内播迁","spread"),("海外播迁","overseas")]:
    names = "、".join(x.get('name','') for x in cities[key])
    lines.append(((f"{label}：{names}"), 14, C_WHITE, False))
add(s, 1.0, 3.0, 11.3, 3.6, lines)

# 长廊概览
s = prs.slides.add_slide(blank); bg(s, "FBF8F2")
add(s, 1.0, 0.7, 11.3, 1.0, [("数字非遗长廊：40 件客家器物", 32, C_INDIGO, True)])
add(s, 1.0, 1.9, 11.3, 0.8, [("瀑布流展陈 · 360° 翻转卡片 · 断代/工艺/流传/谱系/故事", 16, C_EARTH, False)])
names = "、".join(r0['name'] for r0 in relics[:14])
add(s, 1.0, 2.9, 11.3, 2.4, [(names + " ……", 15, C_INK, False)])

# 资源页
s = prs.slides.add_slide(blank); bg(s, INDIGO)
add(s, 1.0, 0.7, 11.3, 1.0, [("学习资源与档案馆", 32, C_EARTH, True)])
lines = [
    ("• 数字门户：首页 / 章节 / 地图 / 长廊 / 档案馆", 16, C_WHITE, False),
    ("• AI 讲解员「阿蛮」：随问随答，流式呈现 + 语音播报", 16, C_WHITE, False),
    ("• 本讲义（Word）与配套 PPT 可于档案馆下载", 16, C_WHITE, False),
    ("• 延展阅读：罗香林《客家研究导论》、各级非遗名录", 16, C_WHITE, False),
]
add(s, 1.0, 2.2, 11.3, 3.6, lines)

prs.save(os.path.join(BASE, "downloads", "meizhou-hakka-heritage-slides.pptx"))
print("✓ pptx 生成完成")
